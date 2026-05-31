#!/usr/bin/env python3
"""Tạo file PowerPoint tổng hợp dự án Nello."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─────────────────────────── helpers ────────────────────────────

BLUE_DARK  = RGBColor(0x1E, 0x3A, 0x5F)   # navy
BLUE_MID   = RGBColor(0x27, 0x6E, 0xB5)   # primary
BLUE_LIGHT = RGBColor(0xD6, 0xE8, 0xF7)   # background accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT  = RGBColor(0x44, 0x44, 0x44)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
RED        = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

IMG_DIR = os.path.join(os.path.dirname(__file__), "src", "images")


def add_bg(slide, color: RGBColor):
    """Fill slide background with solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txb(slide, text, left, top, width, height,
        size=24, bold=False, color=GRAY_TEXT,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def header_bar(slide, title: str, subtitle: str = ""):
    """Dark top bar with title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), BLUE_DARK)
    txb(slide, title,
        Inches(0.3), Inches(0.1),
        Inches(12.0), Inches(0.7),
        size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txb(slide, subtitle,
            Inches(0.3), Inches(0.75),
            Inches(12.0), Inches(0.4),
            size=16, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)


def bullet_list(slide, items, left, top, width, height,
                size=18, color=GRAY_TEXT, title=None, title_color=BLUE_MID):
    if title:
        txb(slide, title, left, top, width, Inches(0.4),
            size=20, bold=True, color=title_color)
        top += Inches(0.45)
        height -= Inches(0.45)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color


def stat_box(slide, label, value, left, top, w, h,
             bg=BLUE_MID, fg=WHITE, label_color=None):
    if label_color is None:
        label_color = RGBColor(0xCC, 0xE0, 0xF5)
    add_rect(slide, left, top, w, h, bg)
    txb(slide, value, left, top + Inches(0.12), w, Inches(0.55),
        size=32, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txb(slide, label, left, top + Inches(0.65), w, Inches(0.35),
        size=13, color=label_color, align=PP_ALIGN.CENTER)


def add_image_safe(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        add_rect(slide, left, top, width, height, RGBColor(0xDD, 0xDD, 0xDD))
        txb(slide, "[Hình ảnh]", left, top + height // 2 - Inches(0.2),
            width, Inches(0.4), size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ─────────────────────────── slides ─────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]   # completely blank


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 – TRANG BÌA
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, BLUE_DARK)
add_rect(sl, 0, Inches(2.5), SLIDE_W, Inches(3.2), BLUE_MID)

txb(sl, "HỆ THỐNG QUẢN LÝ CÔNG VIỆC",
    Inches(0.5), Inches(0.5), Inches(12.0), Inches(0.9),
    size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "NELLO",
    Inches(0.5), Inches(1.3), Inches(12.0), Inches(1.1),
    size=72, bold=True, color=RGBColor(0xFF, 0xD7, 0x00), align=PP_ALIGN.CENTER)

txb(sl, "Ứng dụng quản lý công việc theo mô hình Kanban",
    Inches(1.0), Inches(2.65), Inches(11.0), Inches(0.6),
    size=22, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

txb(sl, "ReactJS  •  Java Spring Boot  •  PostgreSQL  •  AWS",
    Inches(1.0), Inches(3.3), Inches(11.0), Inches(0.5),
    size=18, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(3.5), Inches(4.0), Inches(6.0), Inches(0.05),
         RGBColor(0xFF, 0xD7, 0x00))

txb(sl, "Project I  |  Nguyễn Trọng Nhân  |  MSSV: 202490073",
    Inches(0.5), Inches(4.3), Inches(12.0), Inches(0.5),
    size=16, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
txb(sl, "Thời gian thực hiện: 30/03/2026 – 24/05/2026",
    Inches(0.5), Inches(4.8), Inches(12.0), Inches(0.4),
    size=15, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 – GIỚI THIỆU ĐỀ TÀI
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "GIỚI THIỆU ĐỀ TÀI", "Bài toán và bối cảnh")

# Left panel – bài toán
add_rect(sl, Inches(0.3), Inches(1.4), Inches(5.8), Inches(5.7), BLUE_LIGHT)
txb(sl, "Bài toán",
    Inches(0.5), Inches(1.55), Inches(5.4), Inches(0.4),
    size=20, bold=True, color=BLUE_DARK)
bullet_list(sl, [
    "Số lượng công việc ngày càng tăng cao, khó kiểm soát.",
    "Thông tin bị phân mảnh trên nhiều kênh khác nhau.",
    "Thiếu cái nhìn tổng quan về tiến độ thực hiện dự án.",
    "Suy giảm hiệu suất cộng tác trong nhóm làm việc.",
], Inches(0.5), Inches(2.0), Inches(5.4), Inches(3.5), size=17)

# Right panel – giải pháp
add_rect(sl, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.7), RGBColor(0xE8, 0xF5, 0xE9))
txb(sl, "Giải pháp — Nello",
    Inches(6.7), Inches(1.55), Inches(6.1), Inches(0.4),
    size=20, bold=True, color=RGBColor(0x1B, 0x6C, 0x3A))
bullet_list(sl, [
    "Giao diện Kanban trực quan, quen thuộc.",
    "Phân cấp rõ ràng: Workspace → Board → List → Card.",
    "Kéo thả công việc linh hoạt (Drag & Drop).",
    "Quản lý thành viên, phân quyền theo vai trò.",
    "Tích hợp bảo mật AWS Cognito + HttpOnly Cookie.",
], Inches(6.7), Inches(2.0), Inches(6.1), Inches(3.8), size=17)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 – MỤC TIÊU & PHẠM VI
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "MỤC TIÊU & PHẠM VI ĐỀ TÀI")

# Mục tiêu
bullet_list(sl, [
    "Xây dựng nền tảng quản lý công việc theo mô hình Kanban.",
    "Quản lý thẻ công việc chi tiết (tiêu đề, mô tả, due date, checklist, thành viên).",
    "Triển khai tính năng kéo thả (Drag & Drop) mượt mà.",
    "Đồng bộ dữ liệu nhanh chóng, hỗ trợ cộng tác nhóm.",
],
    Inches(0.4), Inches(1.35), Inches(6.0), Inches(3.2),
    size=17, title="Mục tiêu chính")

# Phạm vi
bullet_list(sl, [
    "Xác thực tài khoản (đăng ký, OTP, đăng nhập, đăng xuất).",
    "Quản lý Workspace, Board, List, Card, Checklist.",
    "Kéo thả Card và List trong Board.",
    "Gán thành viên vào Card; phân quyền ADMIN/MEMBER/VIEWER.",
    "Triển khai trên nền tảng AWS (EC2, RDS, S3, ALB, Cognito, CloudFront).",
],
    Inches(6.8), Inches(1.35), Inches(6.1), Inches(3.6),
    size=17, title="Phạm vi thực hiện")

# Đối tượng
add_rect(sl, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.1), BLUE_LIGHT)
txb(sl, "Đối tượng sử dụng",
    Inches(0.6), Inches(5.1), Inches(12.0), Inches(0.4),
    size=18, bold=True, color=BLUE_DARK)
for i, (icon, text) in enumerate([
    ("👤", "Cá nhân cần sắp xếp công việc cá nhân hàng ngày."),
    ("👥", "Nhóm làm việc / phòng ban cần không gian cộng tác thống nhất."),
    ("📊", "Quản lý dự án cần theo dõi tiến độ và phân bổ nguồn lực."),
]):
    txb(sl, f"{icon}  {text}",
        Inches(0.6 + i * 4.2), Inches(5.55), Inches(4.0), Inches(0.5),
        size=15, color=GRAY_TEXT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 – CÔNG NGHỆ SỬ DỤNG
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "CÔNG NGHỆ SỬ DỤNG", "Stack công nghệ của hệ thống Nello")

tech_groups = [
    ("Frontend", BLUE_MID, [
        "ReactJS — Thư viện UI component-based",
        "Vite — Build tool hiện đại",
        "React DnD — Kéo thả (Drag & Drop)",
        "Axios — HTTP Client",
    ]),
    ("Backend", BLUE_DARK, [
        "Java Spring Boot — RESTful API",
        "Spring Data JPA + Hibernate",
        "Spring Security — Bảo mật",
        "Mô hình MVC phân lớp rõ ràng",
    ]),
    ("Database & Infra", GREEN, [
        "PostgreSQL — CSDL quan hệ",
        "Docker + Nginx — Container hóa",
        "JWT (HttpOnly Cookie) — Xác thực",
    ]),
]

for col, (group, color, items) in enumerate(tech_groups):
    x = Inches(0.3 + col * 4.35)
    w = Inches(4.1)
    add_rect(sl, x, Inches(1.35), w, Inches(0.5), color)
    txb(sl, group, x, Inches(1.37), w, Inches(0.45),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullet_list(sl, items, x, Inches(1.9), w, Inches(4.8), size=16)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 – DỊCH VỤ AWS
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "NỀN TẢNG ĐIỆN TOÁN ĐÁM MÂY AWS", "Các dịch vụ được tích hợp")

services = [
    ("Amazon Cognito", "Quản lý xác thực người dùng, phát hành JWT token."),
    ("Amazon EC2", "Máy chủ ảo triển khai Frontend (Nginx+Docker) và Backend (Spring Boot)."),
    ("Amazon RDS", "Hosting cơ sở dữ liệu PostgreSQL trên đám mây."),
    ("Amazon ALB", "Cân bằng tải (Layer 7), điểm vào duy nhất từ Internet."),
    ("Amazon S3", "Lưu trữ tệp đa phương tiện (avatar, tệp đính kèm)."),
    ("Amazon CloudFront", "CDN phân phối nội dung tĩnh từ S3 đến người dùng."),
]

cols = 3
for i, (name, desc) in enumerate(services):
    row = i // cols
    col = i % cols
    x = Inches(0.3 + col * 4.35)
    y = Inches(1.5 + row * 2.5)
    w = Inches(4.1)
    h = Inches(2.2)
    add_rect(sl, x, y, w, h, BLUE_LIGHT)
    add_rect(sl, x, y, w, Inches(0.45), BLUE_MID)
    txb(sl, name, x, y + Pt(3), w, Inches(0.42),
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, desc, x + Inches(0.1), y + Inches(0.5), w - Inches(0.2), h - Inches(0.6),
        size=14, color=GRAY_TEXT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 – KIẾN TRÚC HỆ THỐNG
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "KIẾN TRÚC HỆ THỐNG", "Triển khai trên AWS Cloud")

arch_img = os.path.join(IMG_DIR, "chap2", "Nello System Architech.png")
add_image_safe(sl, arch_img, Inches(0.3), Inches(1.35), Inches(8.0), Inches(5.7))

bullet_list(sl, [
    "Client truy cập qua ALB (entry point duy nhất).",
    "EC2 Frontend: ReactJS + Nginx + Docker (private subnet).",
    "EC2 Backend: Spring Boot (private subnet).",
    "RDS PostgreSQL: chỉ Backend được truy cập.",
    "Cognito: xác thực, JWT token quản lý bởi Backend.",
    "S3 + CloudFront: phân phối tài nguyên tĩnh.",
],
    Inches(8.5), Inches(1.45), Inches(4.5), Inches(5.5),
    size=15, title="Luồng chính")


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 – PHÂN TÍCH YÊU CẦU CHỨC NĂNG
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "PHÂN TÍCH YÊU CẦU CHỨC NĂNG", "7 nhóm chức năng — 31 Use Case")

func_groups = [
    ("Nhóm 1", "Xác thực & Tài khoản", "7 FR", BLUE_MID),
    ("Nhóm 2", "Quản lý Workspace", "5 FR", BLUE_DARK),
    ("Nhóm 3", "Quản lý Board", "5 FR", GREEN),
    ("Nhóm 4", "Quản lý List", "3 FR", ORANGE),
    ("Nhóm 5", "Quản lý Card", "5 FR", RED),
    ("Nhóm 6", "Kéo thả (DnD)", "2 FR", RGBColor(0x8E, 0x44, 0xAD)),
    ("Nhóm 7", "Quản lý Checklist", "4 FR", RGBColor(0x17, 0x7E, 0x89)),
]

rows, cols_n = 2, 4
for i, (grp, name, count, color) in enumerate(func_groups):
    row = i // cols_n
    col = i % cols_n
    x = Inches(0.3 + col * 3.2)
    y = Inches(1.5 + row * 2.8)
    w = Inches(3.0)
    h = Inches(2.5)
    add_rect(sl, x, y, w, h, RGBColor(0xF5, 0xF5, 0xF5))
    add_rect(sl, x, y, w, Inches(0.15), color)
    txb(sl, grp, x, y + Inches(0.2), w, Inches(0.3),
        size=12, bold=False, color=color, align=PP_ALIGN.CENTER)
    txb(sl, name, x, y + Inches(0.5), w, Inches(0.55),
        size=16, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
    txb(sl, count, x, y + Inches(1.1), w, Inches(0.55),
        size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(sl, "Use Case", x, y + Inches(1.65), w, Inches(0.4),
        size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 – MÔ HÌNH DỮ LIỆU
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "MÔ HÌNH DỮ LIỆU", "Cấu trúc phân cấp 6 thực thể chính")

# Hierarchy boxes
hierarchy = [
    ("User", BLUE_DARK),
    ("Workspace", BLUE_MID),
    ("Board", GREEN),
    ("List", ORANGE),
    ("Card", RED),
    ("Checklist / ChecklistItem", RGBColor(0x8E, 0x44, 0xAD)),
]
box_w = Inches(3.5)
box_h = Inches(0.6)
start_x = Inches(2.5)
start_y = Inches(1.5)
gap = Inches(0.85)

for i, (name, color) in enumerate(hierarchy):
        indent = Inches(i * 0.3)
        x = start_x + indent
        y = start_y + i * gap
        add_rect(sl, x, y, box_w - indent, box_h, color)
        txb(sl, name, x, y + Inches(0.07), box_w - indent, box_h - Inches(0.07),
            size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i > 0:
            # arrow indicator (just text)
            txb(sl, "↑ 1 - N", x - Inches(0.6), y, Inches(0.55), box_h,
                size=11, color=GRAY_TEXT, align=PP_ALIGN.RIGHT)

# Right panel – relation notes
bullet_list(sl, [
    "1 User → nhiều Workspace, Board (thành viên).",
    "1 Workspace → nhiều Board.",
    "1 Board → nhiều List; bảng board_members.",
    "1 List → nhiều Card (sắp xếp theo position).",
    "1 Card → nhiều Checklist; bảng card_members.",
    "1 Checklist → nhiều ChecklistItem.",
    "Xóa CASCADE đảm bảo toàn vẹn dữ liệu.",
    "Position kiểu TEXT hỗ trợ Drag & Drop linh hoạt.",
],
    Inches(7.5), Inches(1.5), Inches(5.4), Inches(5.5),
    size=16, title="Quan hệ chính")


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 – THIẾT KẾ CSDL
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "THIẾT KẾ CƠ SỞ DỮ LIỆU", "PostgreSQL — 10 bảng nghiệp vụ")

db_img = os.path.join(IMG_DIR, "DatabaseDesign.png")
add_image_safe(sl, db_img, Inches(0.3), Inches(1.35), Inches(7.8), Inches(5.7))

bullet_list(sl, [
    "users — UUID PK, liên kết Cognito ID.",
    "workspaces — UUID PK, FK → users.",
    "workspace_members — N-N: users × workspaces + role.",
    "boards — UUID PK, FK → workspaces.",
    "board_members — N-N: users × boards + role.",
    "lists — UUID, position TEXT, FK → boards.",
    "cards — UUID, position TEXT, FK → lists.",
    "card_members — N-N: users × cards.",
    "checklists — UUID, FK → cards.",
    "checklist_items — UUID, is_completed BOOLEAN.",
],
    Inches(8.3), Inches(1.45), Inches(4.7), Inches(5.5),
    size=14, title="Các bảng chính")


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 – LUỒNG XÁC THỰC
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "LUỒNG XÁC THỰC NGƯỜI DÙNG", "Tích hợp Amazon Cognito")

auth_img = os.path.join(IMG_DIR, "chap3", "Auth Flow.png")
add_image_safe(sl, auth_img, Inches(0.3), Inches(1.35), Inches(7.5), Inches(5.7))

steps = [
    "1. Người dùng nhập email + mật khẩu trên ReactJS.",
    "2. Frontend gửi POST /api/v1/auth/sign-in đến Backend qua ALB.",
    "3. Backend gọi Amazon Cognito API xác thực thông tin.",
    "4. Cognito trả về Access Token (JWT) + Refresh Token.",
    "5. Backend đặt token vào HttpOnly Cookie (bảo vệ chống XSS).",
    "6. Frontend chuyển hướng người dùng đến Dashboard.",
    "7. Mỗi API call tiếp theo: Backend xác minh JWT từ Cookie.",
    "8. Token hết hạn → tự động làm mới bằng Refresh Token.",
]

bullet_list(sl, steps,
    Inches(8.0), Inches(1.5), Inches(5.0), Inches(5.5),
    size=15, title="Các bước xác thực")


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 – GIAO DIỆN CHÍNH (1)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "GIAO DIỆN ỨNG DỤNG (1/2)", "Xác thực và Dashboard")

imgs_row1 = [
    (os.path.join(IMG_DIR, "chap4", "screen-signup.png"), "Đăng ký tài khoản"),
    (os.path.join(IMG_DIR, "chap4", "screen-signin.png"), "Đăng nhập"),
    (os.path.join(IMG_DIR, "chap4", "screen-dashboard.png"), "Dashboard"),
]

img_w = Inches(4.0)
img_h = Inches(2.5)
y_img = Inches(1.5)
y_cap = Inches(4.05)

for col, (path, caption) in enumerate(imgs_row1):
    x = Inches(0.3 + col * 4.35)
    add_image_safe(sl, path, x, y_img, img_w, img_h)
    add_rect(sl, x, y_cap, img_w, Inches(0.4), BLUE_MID)
    txb(sl, caption, x, y_cap + Inches(0.05), img_w, Inches(0.35),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom descriptions
descs = [
    "Form nhập fullName, email, password.\nValidation client-side.\nOTP xác nhận qua email.",
    "Form email + password.\nToken lưu HttpOnly Cookie.\nChuyển hướng Dashboard.",
    "Danh sách Workspace / Board.\nThumbnail ảnh nền.\nNút tạo Workspace / Board.",
]
for col, desc in enumerate(descs):
    x = Inches(0.3 + col * 4.35)
    txb(sl, desc, x, Inches(4.5), img_w, Inches(2.5), size=13, color=GRAY_TEXT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 – GIAO DIỆN CHÍNH (2)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "GIAO DIỆN ỨNG DỤNG (2/2)", "Board Kanban — tính năng cốt lõi")

board_img = os.path.join(IMG_DIR, "chap4", "screen-board-kanban.png")
dd_img    = os.path.join(IMG_DIR, "chap4", "feature-drag-drop.png")
card_img  = os.path.join(IMG_DIR, "chap4", "feature-card-detail.png")

add_image_safe(sl, board_img, Inches(0.3), Inches(1.35), Inches(8.5), Inches(3.8))
add_rect(sl, Inches(0.3), Inches(5.2), Inches(8.5), Inches(0.4), BLUE_MID)
txb(sl, "Board Kanban — giao diện chính",
    Inches(0.3), Inches(5.22), Inches(8.5), Inches(0.36),
    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_image_safe(sl, dd_img,   Inches(9.0), Inches(1.35), Inches(4.0), Inches(2.1))
add_image_safe(sl, card_img, Inches(9.0), Inches(3.55), Inches(4.0), Inches(2.1))
txb(sl, "Drag & Drop",
    Inches(9.0), Inches(3.48), Inches(4.0), Inches(0.3),
    size=12, bold=True, color=BLUE_MID)
txb(sl, "Card Detail Modal",
    Inches(9.0), Inches(5.65), Inches(4.0), Inches(0.3),
    size=12, bold=True, color=BLUE_MID)

bullet_list(sl, [
    "Kanban: List dạng cột ngang.",
    "Tạo List / Card inline.",
    "Kéo thả Card & List.",
    "Modal chi tiết Card.",
    "Checklist + Progress Bar.",
    "Due Date badge (cảnh báo).",
    "Avatar thành viên trên Card.",
],
    Inches(9.0), Inches(5.95), Inches(4.0), Inches(1.5),
    size=13, color=GRAY_TEXT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 – KẾT QUẢ ĐẠT ĐƯỢC (THỐNG KÊ)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "KẾT QUẢ ĐẠT ĐƯỢC", "Sau 8 tuần phát triển (30/03 – 24/05/2026)")

stats = [
    ("27/31\nchức năng", "Tỷ lệ hoàn thành", GREEN),
    ("43\nAPI endpoints", "Trên 14 module", BLUE_MID),
    ("6\nmàn hình", "Frontend hoàn chỉnh", BLUE_DARK),
    ("30+\ntest case", "Tỷ lệ pass ~90%", ORANGE),
]
sw = Inches(2.8)
sh = Inches(1.5)
for i, (val, lbl, color) in enumerate(stats):
    x = Inches(0.5 + i * 3.1)
    stat_box(sl, lbl, val, x, Inches(1.45), sw, sh, bg=color)

txb(sl, "Đối chiếu với mục tiêu đề ra:",
    Inches(0.4), Inches(3.2), Inches(12.5), Inches(0.4),
    size=18, bold=True, color=BLUE_DARK)

goals = [
    ("✅", "Xây dựng nền tảng quản lý khoa học",
     "Hỗ trợ đầy đủ Workspace → Board → List → Card."),
    ("✅", "Quản lý thẻ công việc chi tiết",
     "Card đầy đủ: tiêu đề, mô tả, due date, checklist, thành viên."),
    ("✅", "Tính năng kéo thả trực quan",
     "Drag & Drop List và Card ổn định với optimistic update."),
    ("⚠️", "Đồng bộ dữ liệu & cộng tác",
     "Đồng bộ sau mỗi thao tác; chưa có realtime WebSocket."),
]

for i, (icon, title, detail) in enumerate(goals):
    row = i // 2
    col = i % 2
    x = Inches(0.4 + col * 6.5)
    y = Inches(3.7 + row * 1.7)
    add_rect(sl, x, y, Inches(6.2), Inches(1.5),
             BLUE_LIGHT if icon == "✅" else RGBColor(0xFF, 0xF3, 0xE0))
    txb(sl, f"{icon}  {title}", x + Inches(0.1), y + Inches(0.1),
        Inches(6.0), Inches(0.4), size=16, bold=True, color=BLUE_DARK)
    txb(sl, detail, x + Inches(0.3), y + Inches(0.55),
        Inches(5.8), Inches(0.5), size=14, color=GRAY_TEXT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 – HẠN CHẾ & TỒN TẠI
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "CÁC VẤN ĐỀ CÒN TỒN TẠI", "Hạn chế kỹ thuật và chức năng chưa hoàn thiện")

issues = [
    ("Frontend chưa hoàn thiện",
     "Cập nhật/xóa Workspace (FR-08, FR-30) và Board (FR-12, FR-31) — Backend đã xong nhưng UI chưa có.",
     RED),
    ("Lỗi kéo thả vị trí đầu/cuối",
     "Thuật toán beforeId/afterId chưa xử lý edge case khi kéo về đầu hoặc cuối danh sách.",
     ORANGE),
    ("Phân quyền VIEWER chưa đầy đủ",
     "Vai trò VIEWER chưa được kiểm soát tại API write endpoint và chưa ẩn UI tương ứng.",
     ORANGE),
    ("Chưa triển khai lên Staging",
     "Xung đột HttpOnly Cookie cross-domain khi Frontend/Backend trên domain khác nhau.",
     RED),
    ("Chưa có đồng bộ thời gian thực",
     "Request-response thuần túy; thành viên khác cần F5 để thấy cập nhật mới.",
     ORANGE),
    ("Kiểm thử chưa bao phủ edge case",
     "Chủ yếu kiểm thử happy path; thiếu kịch bản lỗi nâng cao và concurrent access.",
     RGBColor(0x7F, 0x8C, 0x8D)),
]

for i, (title, desc, color) in enumerate(issues):
    row = i // 2
    col = i % 2
    x = Inches(0.3 + col * 6.5)
    y = Inches(1.4 + row * 1.95)
    bw = Inches(6.2)
    bh = Inches(1.8)
    add_rect(sl, x, y, bw, bh, RGBColor(0xFD, 0xF0, 0xF0))
    add_rect(sl, x, y, Inches(0.12), bh, color)
    txb(sl, title, x + Inches(0.2), y + Inches(0.1), bw - Inches(0.3), Inches(0.4),
        size=16, bold=True, color=BLUE_DARK)
    txb(sl, desc, x + Inches(0.2), y + Inches(0.55), bw - Inches(0.3), Inches(1.1),
        size=14, color=GRAY_TEXT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 – PHƯƠNG HƯỚNG PHÁT TRIỂN & KẾT LUẬN
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl, WHITE)
header_bar(sl, "PHƯƠNG HƯỚNG PHÁT TRIỂN & KẾT LUẬN")

# Left – hướng phát triển
roadmap = [
    ("🔴 Ưu tiên cao", [
        "Sửa lỗi thuật toán kéo thả edge case.",
        "Hoàn thiện phân quyền VIEWER.",
        "Giải quyết cross-domain cookie → triển khai Staging.",
        "Xây dựng UI cập nhật/xóa Workspace & Board.",
    ]),
    ("🟡 Ưu tiên trung bình", [
        "Realtime WebSocket (Spring STOMP + SockJS).",
        "Quản lý nhãn dán (Tag) cho Card.",
        "Hệ thống thông báo trong ứng dụng.",
    ]),
    ("🟢 Ưu tiên thấp", [
        "Nhật ký hoạt động (Activity Log).",
        "Calendar View / Timeline View.",
        "Upload tệp đính kèm qua Amazon S3.",
        "Caching Redis, phân trang, tối ưu query.",
    ]),
]

y_cur = Inches(1.4)
for prio, items in roadmap:
    txb(sl, prio, Inches(0.3), y_cur, Inches(6.3), Inches(0.35),
        size=16, bold=True, color=BLUE_DARK)
    y_cur += Inches(0.38)
    for item in items:
        txb(sl, "  • " + item, Inches(0.5), y_cur, Inches(6.1), Inches(0.32),
            size=14, color=GRAY_TEXT)
        y_cur += Inches(0.32)
    y_cur += Inches(0.1)

# Right – kết luận
add_rect(sl, Inches(7.0), Inches(1.4), Inches(6.0), Inches(5.7), BLUE_DARK)
txb(sl, "Kết luận",
    Inches(7.2), Inches(1.55), Inches(5.6), Inches(0.5),
    size=22, bold=True, color=WHITE)
bullet_list(sl, [
    "Đã xây dựng thành công hệ thống quản lý công việc Nello theo mô hình Kanban.",
    "Hoàn thành 27/31 chức năng (87,1%), 43 API endpoint, 6 màn hình Frontend.",
    "Kiến trúc Frontend/Backend tách biệt, tích hợp AWS Cognito đảm bảo bảo mật.",
    "Trải nghiệm UX tốt với Drag & Drop và optimistic update.",
    "Hệ thống chứng minh tính khả thi và sẵn sàng mở rộng.",
],
    Inches(7.2), Inches(2.2), Inches(5.6), Inches(4.6),
    size=15, color=WHITE)


# ─────────────────────────── save ────────────────────────────
out = os.path.join(os.path.dirname(__file__), "Nello_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
